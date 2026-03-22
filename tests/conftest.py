"""Shared test fixtures."""
from __future__ import annotations

from pathlib import Path

import pikepdf
import pytest
from PIL import Image, ImageDraw
from pikepdf import Array, Dictionary, Name, String
from reportlab.pdfgen import canvas


@pytest.fixture()
def sample_pdf(tmp_path: Path) -> Path:
    """Create a simple 5-page PDF for testing."""
    pdf_path = tmp_path / "sample.pdf"
    pdf = pikepdf.Pdf.new()
    for i in range(5):
        page = pikepdf.Page(pikepdf.Dictionary(
            Type=pikepdf.Name.Page,
            MediaBox=[0, 0, 612, 792],
            Contents=pdf.make_stream(f"BT /F1 12 Tf 100 700 Td (Page {i + 1}) Tj ET".encode()),
            Resources=pikepdf.Dictionary(
                Font=pikepdf.Dictionary(
                    F1=pikepdf.Dictionary(
                        Type=pikepdf.Name.Font,
                        Subtype=pikepdf.Name.Type1,
                        BaseFont=pikepdf.Name.Helvetica,
                    )
                )
            ),
        ))
        pdf.pages.append(page)
    pdf.save(pdf_path)
    return pdf_path


@pytest.fixture()
def encrypted_pdf(tmp_path: Path, sample_pdf: Path) -> Path:
    """Create a password-protected PDF."""
    enc_path = tmp_path / "encrypted.pdf"
    with pikepdf.open(sample_pdf) as pdf:
        pdf.save(
            enc_path,
            encryption=pikepdf.Encryption(user="userpass", owner="ownerpass", R=6),
        )
    return enc_path


@pytest.fixture()
def sample_images(tmp_path: Path) -> list[Path]:
    """Create sample PNG images for testing."""
    paths = []
    for i in range(3):
        img_path = tmp_path / f"img_{i + 1}.png"
        img = Image.new("RGB", (200, 300), color=(50 * i, 100, 200))
        img.save(img_path)
        paths.append(img_path)
    return paths


@pytest.fixture()
def workdir(tmp_path: Path) -> Path:
    """Provide a clean working directory."""
    wd = tmp_path / "workdir"
    wd.mkdir()
    return wd


@pytest.fixture()
def blank_mix_pdf(tmp_path: Path) -> Path:
    """Create a PDF containing both normal and blank pages."""
    target = tmp_path / "blank_mix.pdf"
    c = canvas.Canvas(str(target), pagesize=(595, 842))
    y = 800
    for i in range(40):
        c.drawString(40, y, f"This is a non blank page line {i} with enough text.")
        y -= 18
    c.showPage()
    c.showPage()
    y = 800
    for i in range(40):
        c.drawString(40, y, f"This is another non blank page line {i} with enough text.")
        y -= 18
    c.showPage()
    c.showPage()
    c.save()
    return target


@pytest.fixture()
def image_pdf(tmp_path: Path, sample_images: list[Path]) -> Path:
    """Create a PDF that embeds a raster image."""
    target = tmp_path / "image_pdf.pdf"
    c = canvas.Canvas(str(target), pagesize=(300, 300))
    c.drawImage(str(sample_images[0]), 40, 60, width=180, height=180)
    c.showPage()
    c.save()
    return target


@pytest.fixture()
def attachment_pdf(tmp_path: Path, sample_pdf: Path) -> Path:
    """Create a PDF with an embedded attachment."""
    target = tmp_path / "attachment.pdf"
    target.write_bytes(sample_pdf.read_bytes())
    with pikepdf.open(target, allow_overwriting_input=True) as pdf:
        pdf.attachments["note.txt"] = b"attached hello"
        pdf.save(target)
    return target


@pytest.fixture()
def form_pdf(tmp_path: Path) -> Path:
    """Create a PDF with a simple AcroForm text field."""
    target = tmp_path / "form.pdf"
    c = canvas.Canvas(str(target), pagesize=(300, 300))
    form = c.acroForm
    form.textfield(name="name", x=60, y=180, width=160, height=24, value="")
    c.showPage()
    c.save()
    return target


@pytest.fixture()
def signature_field_pdf(tmp_path: Path) -> Path:
    """Create a PDF containing a minimal signature field."""
    target = tmp_path / "signature_field.pdf"
    pdf = pikepdf.Pdf.new()
    pdf.add_blank_page(page_size=(200, 200))
    page = pdf.pages[0]
    sig_dict = pdf.make_indirect(
        Dictionary(
            Type=Name("/Sig"),
            Filter=Name("/Adobe.PPKLite"),
            SubFilter=Name("/adbe.pkcs7.detached"),
            Name=String("Tester"),
        )
    )
    field = pdf.make_indirect(
        Dictionary(
            FT=Name("/Sig"),
            T=String("Signature1"),
            V=sig_dict,
            Rect=Array([0, 0, 0, 0]),
            P=page.obj,
        )
    )
    pdf.Root.AcroForm = Dictionary(Fields=Array([field]), SigFlags=3)
    pdf.save(target)
    return target


@pytest.fixture()
def scanned_pdf(tmp_path: Path) -> Path:
    """Create a scanned-like PDF from an image with text."""
    image_path = tmp_path / "scanned.png"
    pdf_path = tmp_path / "scanned.pdf"
    image = Image.new("RGB", (900, 500), color="white")
    draw = ImageDraw.Draw(image)
    draw.text((80, 200), "Scanned OCR sample", fill="black")
    image.save(image_path)
    image.save(pdf_path, "PDF", resolution=144.0)
    return pdf_path


@pytest.fixture()
def sample_docx(tmp_path: Path) -> Path:
    """Create a simple DOCX file."""
    from docx import Document

    target = tmp_path / "sample.docx"
    doc = Document()
    doc.add_heading("Sample Document", level=1)
    doc.add_paragraph("Hello from DOCX.")
    doc.save(target)
    return target


@pytest.fixture()
def sample_xlsx(tmp_path: Path) -> Path:
    """Create a simple XLSX file."""
    from openpyxl import Workbook

    target = tmp_path / "sample.xlsx"
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Name"
    ws["B1"] = "Value"
    ws["A2"] = "Alpha"
    ws["B2"] = 42
    wb.save(target)
    return target


@pytest.fixture()
def sample_pptx(tmp_path: Path) -> Path:
    """Create a simple PPTX file."""
    from pptx import Presentation

    target = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Sample PPTX"
    slide.placeholders[1].text = "Hello from PowerPoint."
    prs.save(target)
    return target


@pytest.fixture()
def rendered_text_pdf(tmp_path: Path) -> Path:
    """Create a standard text PDF rendered by reportlab."""
    target = tmp_path / "rendered_text.pdf"
    c = canvas.Canvas(str(target), pagesize=(300, 300))
    c.drawString(50, 250, "Rendered page 1")
    c.showPage()
    c.drawString(50, 250, "Rendered page 2")
    c.showPage()
    c.save()
    return target


@pytest.fixture()
def rotated_text_pdf(tmp_path: Path) -> Path:
    """Create a PDF whose page content is rotated 90 degrees with enough text for OSD."""
    target = tmp_path / "rotated_text.pdf"
    c = canvas.Canvas(str(target), pagesize=(595, 842))
    c.translate(300, 420)
    c.rotate(90)
    y = 200
    for i in range(25):
        c.drawString(-250, y, f"This is OCR orientation sample line number {i} with enough text for detection.")
        y -= 20
    c.showPage()
    c.save()
    return target
