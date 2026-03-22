"""Office to PDF conversion via LibreOffice headless mode."""
from __future__ import annotations

import shutil
import textwrap
from pathlib import Path

from pdf_agent.core import ErrorCode, ToolError
from pdf_agent.schemas.tool import ToolInputSpec, ToolManifest, ToolOutputSpec
from pdf_agent.tools.base import BaseTool, ProgressReporter, ToolResult
from pdf_agent.tools.filenames import localized_output_name
from pdf_agent.tools.libreoffice import run_libreoffice_conversion


class OfficeToPdfTool(BaseTool):
    def manifest(self) -> ToolManifest:
        return ToolManifest(
            name="office_to_pdf",
            label="Office 转 PDF",
            category="convert",
            description="将 Word、Excel、PowerPoint 等 Office 文件转为 PDF",
            inputs=ToolInputSpec(
                min=1,
                max=1,
                accept=[
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    "application/msword",
                    "application/vnd.ms-excel",
                    "application/vnd.ms-powerpoint",
                ],
            ),
            outputs=ToolOutputSpec(type="pdf"),
            params=[],
            engine="libreoffice+python-fallback",
            async_hint=True,
        )

    def validate(self, params: dict) -> dict:
        return {}

    def run(self, inputs: list[Path], params: dict, workdir: Path, reporter: ProgressReporter | None = None) -> ToolResult:
        lo_bin = shutil.which("libreoffice") or shutil.which("soffice")
        engine = "python-fallback"
        fallback_used = True
        fallback_reason: str | None = None
        if reporter:
            reporter(10, "Starting Office to PDF conversion")
        default_output_path = workdir / f"{inputs[0].stem}.pdf"
        if lo_bin:
            success, failure_reason = run_libreoffice_conversion(
                lo_bin,
                convert_to="pdf",
                input_path=inputs[0],
                outdir=workdir,
                profile_dir=workdir / ".libreoffice-profile",
            )
            if success and default_output_path.exists():
                engine = "libreoffice"
                fallback_used = False
            else:
                fallback_reason = failure_reason or "LibreOffice did not produce a PDF"
        else:
            fallback_reason = "LibreOffice is not installed"

        if fallback_used:
            self._fallback_convert(inputs[0], default_output_path)

        if not default_output_path.exists():
            raise ToolError(ErrorCode.OUTPUT_GENERATION_FAILED, "Failed to generate PDF output")
        output_path = workdir / localized_output_name(inputs[0], "转PDF")
        if output_path != default_output_path:
            default_output_path.replace(output_path)
        if reporter:
            reporter(100, "Done")
        meta = {"output": output_path.name, "engine": engine, "fallback_used": fallback_used}
        if fallback_reason:
            meta["fallback_reason"] = fallback_reason
        log = f"Converted {inputs[0].name} to PDF via {engine}"
        if fallback_reason:
            log += f" (fallback reason: {fallback_reason})"
        return ToolResult(
            output_files=[output_path],
            meta=meta,
            log=log,
        )

    @staticmethod
    def _fallback_convert(input_path: Path, output_path: Path) -> None:
        suffix = input_path.suffix.lower()
        if suffix == ".docx":
            _docx_to_pdf_fallback(input_path, output_path)
            return
        if suffix == ".xlsx":
            _xlsx_to_pdf_fallback(input_path, output_path)
            return
        if suffix == ".pptx":
            _pptx_to_pdf_fallback(input_path, output_path)
            return
        raise ToolError(
            ErrorCode.OUTPUT_GENERATION_FAILED,
            f"LibreOffice unavailable and no fallback converter for {input_path.suffix or 'this file type'}",
        )


def _docx_to_pdf_fallback(input_path: Path, output_path: Path) -> None:
    try:
        from docx import Document
    except ImportError as exc:
        raise ToolError(ErrorCode.ENGINE_NOT_INSTALLED, "python-docx is not installed") from exc

    document = Document(input_path)
    lines = [input_path.stem]
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                lines.append(" | ".join(values))
    _render_text_lines_to_pdf(lines, output_path)


def _xlsx_to_pdf_fallback(input_path: Path, output_path: Path) -> None:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise ToolError(ErrorCode.ENGINE_NOT_INSTALLED, "openpyxl is not installed") from exc

    workbook = load_workbook(input_path, data_only=True)
    lines = [input_path.stem]
    for worksheet in workbook.worksheets:
        lines.append(f"[Sheet] {worksheet.title}")
        for row in worksheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value not in (None, "")]
            if values:
                lines.append(" | ".join(values))
    _render_text_lines_to_pdf(lines, output_path)


def _pptx_to_pdf_fallback(input_path: Path, output_path: Path) -> None:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ToolError(ErrorCode.ENGINE_NOT_INSTALLED, "python-pptx is not installed") from exc

    presentation = Presentation(input_path)
    lines = [input_path.stem]
    for slide_index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"[Slide {slide_index}]")
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                lines.extend(part.strip() for part in text.splitlines() if part.strip())
    _render_text_lines_to_pdf(lines, output_path)


def _render_text_lines_to_pdf(lines: list[str], output_path: Path) -> None:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
    except ImportError as exc:
        raise ToolError(ErrorCode.ENGINE_NOT_INSTALLED, "reportlab is not installed") from exc

    pdf = canvas.Canvas(str(output_path), pagesize=A4)
    width, height = A4
    x = 48
    y = height - 56
    max_chars = 70
    line_height = 18

    def ensure_space() -> None:
        nonlocal y
        if y < 56:
            pdf.showPage()
            y = height - 56

    for raw_line in lines or [output_path.stem]:
        line = raw_line.strip() or " "
        wrapped = textwrap.wrap(line, width=max_chars) or [" "]
        for part in wrapped:
            ensure_space()
            pdf.drawString(x, y, part)
            y -= line_height
        y -= 4
    pdf.save()
