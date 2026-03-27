"""PDF to Excel/PPT tool — convert PDF to xlsx or pptx via LibreOffice."""
from __future__ import annotations

import shutil
from pathlib import Path

import pikepdf

from pdf_agent.core import ErrorCode, ToolError
from pdf_agent.schemas.tool import ToolInputSpec, ToolManifest, ToolOutputSpec
from pdf_agent.tools._builtins.pdf_to_text import _extract_page_text
from pdf_agent.tools.base import BaseTool, ProgressReporter, ToolResult
from pdf_agent.tools.filenames import localized_output_name
from pdf_agent.tools.libreoffice import run_libreoffice_conversion


class PdfToExcelTool(BaseTool):
    def manifest(self) -> ToolManifest:
        return ToolManifest(
            name="pdf_to_excel",
            label="PDF 转 Excel",
            category="convert",
            description="将 PDF 转换为 Excel 表格文件（.xlsx）",
            inputs=ToolInputSpec(min=1, max=1),
            outputs=ToolOutputSpec(type="xlsx"),
            params=[],
            engine="libreoffice+openpyxl",
            async_hint=True,
        )

    def validate(self, params: dict) -> dict:
        return {}

    def run(self, inputs: list[Path], params: dict, workdir: Path, reporter: ProgressReporter | None = None) -> ToolResult:
        output_path = workdir / localized_output_name(inputs[0], "转Excel", ext=".xlsx")
        lo_bin = shutil.which("libreoffice") or shutil.which("soffice")
        engine = "openpyxl"
        fallback_used = not bool(lo_bin)
        fallback_reason: str | None = "libreoffice not found" if not lo_bin else None
        if lo_bin:
            if reporter:
                reporter(10, "Starting conversion...")
            success, failure_reason = run_libreoffice_conversion(
                lo_bin,
                convert_to="xlsx",
                input_path=inputs[0],
                outdir=workdir,
                profile_dir=workdir / ".libreoffice-profile",
            )
            if success and output_path.exists():
                engine = "libreoffice"
            else:
                fallback_used = True
                fallback_reason = failure_reason or "LibreOffice did not produce a .xlsx file"
        if engine != "libreoffice":
            self._fallback_convert(inputs[0], output_path)
        if reporter:
            reporter(100, "Done")
        meta = {"output": output_path.name, "engine": engine, "fallback_used": fallback_used}
        if fallback_reason:
            meta["fallback_reason"] = fallback_reason
        log = f"Converted to {output_path.name} via {engine}"
        if fallback_reason:
            log += f" (fallback reason: {fallback_reason})"
        return ToolResult(output_files=[output_path], meta=meta, log=log)

    @staticmethod
    def _fallback_convert(input_pdf: Path, output_path: Path) -> None:
        try:
            from openpyxl import Workbook
        except ImportError as exc:
            raise ToolError(ErrorCode.ENGINE_NOT_INSTALLED, "openpyxl is not installed") from exc

        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "PDF Text"
        sheet.append(["Page", "Line", "Text"])
        with pikepdf.open(input_pdf) as pdf:
            for page_index, page in enumerate(pdf.pages, start=1):
                lines = [line.strip() for line in _extract_page_text(page).splitlines() if line.strip()]
                if not lines:
                    lines = ["[No extractable text]"]
                for line_index, line in enumerate(lines, start=1):
                    sheet.append([page_index, line_index, line])
        workbook.save(output_path)


class PdfToPptTool(BaseTool):
    def manifest(self) -> ToolManifest:
        return ToolManifest(
            name="pdf_to_ppt",
            label="PDF 转 PPT",
            category="convert",
            description="将 PDF 转换为 PowerPoint 演示文稿（.pptx）",
            inputs=ToolInputSpec(min=1, max=1),
            outputs=ToolOutputSpec(type="pptx"),
            params=[],
            engine="libreoffice+python-pptx",
            async_hint=True,
        )

    def validate(self, params: dict) -> dict:
        return {}

    def run(self, inputs: list[Path], params: dict, workdir: Path, reporter: ProgressReporter | None = None) -> ToolResult:
        output_path = workdir / localized_output_name(inputs[0], "转PPT", ext=".pptx")
        lo_bin = shutil.which("libreoffice") or shutil.which("soffice")
        engine = "python-pptx"
        fallback_used = not bool(lo_bin)
        fallback_reason: str | None = "libreoffice not found" if not lo_bin else None
        if lo_bin:
            if reporter:
                reporter(10, "Starting conversion...")
            success, failure_reason = run_libreoffice_conversion(
                lo_bin,
                convert_to="pptx",
                input_path=inputs[0],
                outdir=workdir,
                profile_dir=workdir / ".libreoffice-profile",
            )
            if success and output_path.exists():
                engine = "libreoffice"
            else:
                fallback_used = True
                fallback_reason = failure_reason or "LibreOffice did not produce a .pptx file"
        if engine != "libreoffice":
            self._fallback_convert(inputs[0], output_path)
        if reporter:
            reporter(100, "Done")
        meta = {"output": output_path.name, "engine": engine, "fallback_used": fallback_used}
        if fallback_reason:
            meta["fallback_reason"] = fallback_reason
        log = f"Converted to {output_path.name} via {engine}"
        if fallback_reason:
            log += f" (fallback reason: {fallback_reason})"
        return ToolResult(output_files=[output_path], meta=meta, log=log)

    @staticmethod
    def _fallback_convert(input_pdf: Path, output_path: Path) -> None:
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError as exc:
            raise ToolError(ErrorCode.ENGINE_NOT_INSTALLED, "python-pptx is not installed") from exc

        presentation = Presentation()
        with pikepdf.open(input_pdf) as pdf:
            for page_index, page in enumerate(pdf.pages, start=1):
                slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                slide.shapes.title.text = f"Page {page_index}"
                text = _extract_page_text(page).strip() or "[No extractable text]"
                body = slide.placeholders[1]
                body.text = text
                body.width = Inches(8.5)
                body.height = Inches(5.0)
        presentation.save(output_path)
